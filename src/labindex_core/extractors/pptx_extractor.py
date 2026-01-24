"""
PowerPoint document extractor.

Handles .pptx files using python-pptx.
"""

from pathlib import Path
from typing import List

from .base import TextExtractor, ExtractionResult


class PowerPointExtractor(TextExtractor):
    """Extract text from PowerPoint presentations."""

    EXTENSIONS = ['.pptx', '.ppt']

    def _extract_impl(self, path: Path) -> ExtractionResult:
        """Extract text from PowerPoint."""
        ext = path.suffix.lower()

        if ext == '.pptx':
            return self._extract_pptx(path)
        elif ext == '.ppt':
            return ExtractionResult.failure(".ppt format not supported (use .pptx)")
        else:
            return ExtractionResult.failure(f"Unknown PowerPoint format: {ext}")

    def _extract_pptx(self, path: Path) -> ExtractionResult:
        """Extract from .pptx using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            return ExtractionResult.failure("python-pptx not installed")

        try:
            prs = Presentation(path)

            all_text: List[str] = []
            sources = {}

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text: List[str] = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)

                    # Handle tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(' | '.join(row_text))

                if slide_text:
                    slide_content = '\n'.join(slide_text)
                    all_text.append(f"[Slide {slide_num}]\n{slide_content}")
                    sources[f"slide_{slide_num}"] = slide_content[:500]

            return ExtractionResult(
                text='\n\n'.join(all_text),
                metadata={'slide_count': len(prs.slides)},
                sources=sources
            )

        except Exception as e:
            return ExtractionResult.failure(f"PowerPoint read error: {e}")
